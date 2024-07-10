import base64
import concurrent.futures
import logging
import os
import tempfile
import time

import streamlit as st
from PIL import Image
from langchain.chains import TransformChain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate, HumanMessagePromptTemplate
from langchain_core.prompts import PromptTemplate
from langchain_core.pydantic_v1 import BaseModel, Field
from langchain_core.runnables import RunnablePassthrough
from langchain_core.runnables import chain
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from pptx import Presentation
import aspose.slides as slides
import aspose.pydrawing as drawing

# 로깅 설정
logging.basicConfig(level=logging.DEBUG)

# 스레드 풀 실행자 초기화
executor = concurrent.futures.ThreadPoolExecutor()


# Set verbose if needed
# globals.set_debug(True)


# # # # # # # # # # # #
# 테스트 코드. #
# # # # # # # # # # # #


# # # # # # # # # # # # # # # # # # # # # # # #
# ppt 파일 전체를 png 이미지로 변환하는 코드 #
# # # # # # # # # # # # # # # # # # # # # # # #


def pptx_to_png(pptx_filename: str, output_folder='tmp_ppt_images_folder'):
    # 이미지 임시 저장 디렉토리 생성
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    logging.debug('PPT 파일 PNG로 변환 시작 :' + pptx_filename)

    png_file_list = []
    with slides.Presentation(pptx_filename) as presentation:
        for slide in presentation.slides:
            # 파일 저장 경로 및 이름.
            tmp_name = os.path.join(output_folder,
                                    "presentation_slide_{0}.png".format(str(slide.slide_number)))
            tmp_name = os.path.abspath(tmp_name)
            logging.debug('저장시작 - ' + tmp_name)

            # png 파일 저장.
            slide.get_thumbnail(1, 1).save(tmp_name, drawing.imaging.ImageFormat.png)
            # png 파일 목록 저장
            png_file_list.append(tmp_name)
            # time.sleep(0.01)
            logging.debug('저장완료 - ' + tmp_name)

    return png_file_list


# # # # # # # # # # # #
# RAG를 위한 체인 #
# # # # # # # # # # # #


# ppt 텍스트 추출 함수(페이지별)
def get_text_documents_from_pptx(uploaded_pptx_file):
    # print('\n'*10, uploaded_pptx_file.getvalue(), '\n'*5)

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(uploaded_pptx_file.getvalue())
        tmp_file_path = tmp_file.name
        tmp_file.close()

        # UnstructuredPowerPointLoader를 사용하여 파일 내용 읽기
        loader = UnstructuredPowerPointLoader(
            tmp_file_path, mode='paged', strategy='fast',
        )
    return loader.load(), tmp_file_path  # document와 임시저장한 파일경로도 리턴.


# 이미지/텍스트 도큐먼트를 합쳐주는 함수
def combine_and_sort_documents(docs_list1: list, docs_list2: list, by='page_number') -> list:
    all_docs = docs_list1 + docs_list2
    all_docs.sort(key=lambda x: x.metadata[by])
    return all_docs


def format_docs(docs):
    # 검색한 문서 결과를 하나의 문단으로 합쳐줍니다.
    return "\n\n".join(doc.page_content for doc in docs)


def get_rag_chain_from_docs(docs_or_vectorstore_for_rag, from_docs=True):
    # text_splitter = RecursiveCharacterTextSplitter(
    #     chunk_size=2000,
    #     chunk_overlap=200
    # )
    # splits = text_splitter.split_documents(docs_for_rag)

    if from_docs:
        # 벡터스토어를 생성
        vectorstore = FAISS.from_documents(
            documents=docs_or_vectorstore_for_rag, embedding=OpenAIEmbeddings())
    else:
        vectorstore = docs_or_vectorstore_for_rag

    # 검색기 생성.
    retriever = vectorstore.as_retriever(k=10)

    # 기본 프롬프트 생성
    rag_default_prompt = ("You are an assistant for question-answering tasks. "
                          "Use the following pieces of retrieved context to answer the question. "
                          "If you don't know the answer, just say that '잘 모르겠어요.' "
                          "Use three sentences maximum and keep the answer concise."
                          "\n"
                          "Question: {question} "
                          "\n"
                          "Context: {context} "
                          "\n"
                          "Answer:")

    my_prompt = ChatPromptTemplate(input_variables=['context', 'question'],
                                   messages=[HumanMessagePromptTemplate(
                                       prompt=PromptTemplate(input_variables=['context', 'question'],
                                                             template=rag_default_prompt))])

    llm = ChatOpenAI(model_name="gpt-4o", temperature=0)

    # rag 체인을 생성합니다.
    rag_chain = (
            {"context": retriever | format_docs, "question": RunnablePassthrough()}
            | my_prompt
            | llm
            | StrOutputParser()
    )

    return rag_chain, vectorstore


# # # # # # # # # # # #
# 이미지 해석 체인 #
# # # # # # # # # # # #


# 이미지 추출/디스크립션 함수
def extract_image_description_from_shape(shape, image_output_dir, slide_number, image_count) -> (dict, str):
    extracted_image = shape.image
    image_bytes = extracted_image.blob

    image_format = extracted_image.ext.lower()
    if image_format != 'jpg' and image_format != 'jpeg' and \
            image_format != 'png':
        return None, None

    # 임시 경로에 이미지를 저장하고 절대 경로로 변환.
    image_filename = os.path.join(image_output_dir, f"slide_{slide_number:03d}_image_{image_count:03d}.{image_format}")
    image_filename = os.path.abspath(image_filename)

    # 이미지 파일로 저장
    with open(image_filename, "wb") as f:
        f.write(image_bytes)

    # 텍스트 디스크립션 생성
    text_desc_gen_result = get_image_description(image_filename)
    # image_description, node_list, link_list

    # 망구성도가 아닐 경우 None 리턴
    if text_desc_gen_result is None or \
            text_desc_gen_result['image_description'] is None or text_desc_gen_result['image_description'] == '' or \
            text_desc_gen_result['node_list'] is None or len(text_desc_gen_result['node_list']) == 0 or \
            text_desc_gen_result['link_list'] is None or len(text_desc_gen_result['link_list']) == 0 or \
            text_desc_gen_result['mermaid_code'] is None or text_desc_gen_result['mermaid_code'] == '':
        return None, None

    return text_desc_gen_result, image_filename


# 이미지에 LLM 디스크립션을 달아주는 함수. (shape 여러개로 쪼개진 이미지들 때문에 미사용)
def get_image_documents_from_pptx_old(pptx_path, image_output_dir='./tmp_extracted_images',
                                      begin_of_image_token='<image>', end_of_image_token='</image>'):
    # 프레젠테이션 열기
    prs = Presentation(pptx_path)

    # 이미지 임시 저장 디렉토리 생성
    if not os.path.exists(image_output_dir):
        os.makedirs(image_output_dir)

    # 메타데이터 생성 - 디렉토리와 파일명
    directory = os.path.dirname(pptx_path.name)
    filename = os.path.basename(pptx_path.name)

    # 도큐먼트를 저장할 리스트 초기화
    docs_list = []
    img_filepath_list = []

    # 슬라이드마다 반복
    for slide_number, slide in enumerate(prs.slides):
        image_count = 0
        # 슬라이드 내의 모든 쉐이프 확인
        for shape in slide.shapes:
            # 이미지 우선 추출.
            if hasattr(shape, "image"):
                img_desc_result, img_filename = extract_image_description_from_shape(shape, image_output_dir,
                                                                                     slide_number, image_count)

                if img_desc_result is None:  # jpg, png 파일 그림이 아니거나 망구성도가 아닐 경우.
                    continue

                img_desc_str = ''
                for k, v in img_desc_result.items():
                    img_desc_str += str(k) + ':\n' + str(v) + '\n'
                img_desc_str = img_desc_str.strip()

                # 특수 토큰 붙이기 (디스크립션 부분만 사용)
                if begin_of_image_token:
                    img_desc_str = begin_of_image_token + '\n' + img_desc_str
                if end_of_image_token:
                    img_desc_str = img_desc_str + '\n' + end_of_image_token

                # 텍스트 생성 결과 저장
                tmp_img_doc = Document(page_content=img_desc_str,
                                       metadata={'source': pptx_path.name,
                                                 'file_directory': directory,
                                                 'filename': filename,
                                                 'page_number': slide_number + 1})

                docs_list.append(tmp_img_doc)
                img_filepath_list.append(img_filename)

    # 이미지 추출 도큐먼트 리스트, 이미지 파일명 리스트를 반환
    return docs_list, img_filepath_list


# pptx 파일의 전체 슬라이드를 모두 png 파일로 저장한 후, 디스크립션을 받아오는 함수.
def get_image_documents_from_pptx(pptx_path: str, image_output_dir='tmp_ppt_images_folder',
                                  begin_of_image_token='<image>', end_of_image_token='</image>'):
    # 프레젠테이션 열기
    image_file_list = pptx_to_png(pptx_path, output_folder=image_output_dir)

    # 메타데이터 생성 - 디렉토리와 파일명
    directory = os.path.dirname(pptx_path)
    filename = os.path.basename(pptx_path)

    # 도큐먼트를 저장할 리스트 초기화
    docs_list = []
    img_filepath_list = []

    # 각 파일마다 디스크립션을 가져온다.
    for slide_number, image_file in enumerate(image_file_list):
        logging.debug('\n슬라이드 이미지 읽는중 : ' + str(slide_number+1) + ' / ' + str(len(image_file_list)))

        # 텍스트 디스크립션 생성
        desc_gen_result = get_image_description(image_file)
        # image_description, node_list, link_list

        # 망구성도가 아닐 경우 skip
        if desc_gen_result is None or \
                desc_gen_result['image_description'] is None or desc_gen_result['image_description'] == '' or \
                desc_gen_result['node_list'] is None or len(desc_gen_result['node_list']) == 0 or \
                desc_gen_result['link_list'] is None or len(desc_gen_result['link_list']) == 0 or \
                desc_gen_result['mermaid_code'] is None or desc_gen_result['mermaid_code'] == '':
            continue

        img_desc_str = ''
        for k, v in desc_gen_result.items():
            img_desc_str += str(k) + ':\n' + str(v) + '\n'
        img_desc_str = img_desc_str.strip()

        # 특수 토큰 붙이기
        if begin_of_image_token:
            img_desc_str = begin_of_image_token + '\n' + img_desc_str
        if end_of_image_token:
            img_desc_str = img_desc_str + '\n' + end_of_image_token

        # 텍스트 생성 결과 저장
        tmp_img_doc = Document(page_content=img_desc_str,
                               metadata={'source': pptx_path,
                                         'file_directory': directory,
                                         'filename': filename,
                                         'page_number': slide_number + 1})

        docs_list.append(tmp_img_doc)
        img_filepath_list.append(image_file)

    # 이미지 추출 도큐먼트 리스트, 이미지 파일명 리스트를 반환
    return docs_list, img_filepath_list


# # # # # # # # # #
# 이미지 체인 생성 #
# # # # # # # # # #


# 이미지 파일을 읽어 인코딩해주는 함수
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')


# 로컬 디스크의 이미지를 불러오는 함수
def load_image_from_filename(inputs: dict) -> dict:
    """Load image from file and encode it as base64."""
    # 파일명일 경우는 읽어서 인코딩 후 리턴.
    image_path = inputs["image_path"]
    image_base64 = encode_image(image_path)
    return {"image": image_base64}


def load_image_from_bytes(inputs: dict) -> dict:
    """Load image from file and encode it as base64."""
    # 이미지 bytes 타입을 utf-8로 디코딩
    img_base64 = base64.b64encode(inputs["image_path"]).decode('utf-8')
    return {"image": img_base64}


# 이미지 정보 클래스
class ImageInformation(BaseModel):
    """Information about an image."""
    image_description: str = Field(description="망 구성에 대한 자세한 디스크립션(한국어)")
    node_list: list[str] = Field(description="망구성도 속의 노드 목록")
    link_list: list[str] = Field(description="망구성도 속의 링크/에지/인터페이스 목록")
    mermaid_code: str = Field(description="Mermaid chart code")


# 체인을 합쳐서 invoke 해주는 함수
def get_image_description(image_path, is_bytes=False) -> dict:
    # 체인1-1 - 파일에서 가져오기
    load_image_chain = TransformChain(
        input_variables=["image_path"],
        output_variables=["image"],
        transform=load_image_from_filename
    )

    # 체인1-2 - 바이츠 객체에서 가져오기
    pass_image_chain = TransformChain(
        input_variables=["image_path"],
        output_variables=["image"],
        transform=load_image_from_bytes
    )

    # 체인3
    parser = JsonOutputParser(pydantic_object=ImageInformation)

    # 체인2
    @chain
    def image_model(inputs: dict) -> str | list[str] | dict:
        """Invoke model with image and prompt."""
        model = ChatOpenAI(temperature=0, model="gpt-4o", max_tokens=4096)  # model="gpt-4-vision-preview"
        msg = model.invoke(
            [HumanMessage(
                content=[
                    {"type": "text", "text": inputs["prompt"]},
                    {"type": "text", "text": parser.get_format_instructions()},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{inputs['image']}"}},
                ])],

        )
        return msg.content

    vision_prompt = """
    Given the image, provide the following information in Korean:
    - 망구성도 이미지에 대한 자세한 디스크립션 (이미지 속 노드들과 노드 간의 연결에 대해 최대한 상세하게 묘사할 것.)
    - 망구성도 이미지 속 모든 노드 이름 목록
    - 망구성도 이미지 속 모든 링크/에지/인터페이스 이름 목록
    - 망구성도를 Mermaid code로 변환한 텍스트 
    단, 이미지가 노드와 링크로 구성된 망구성도 그림이 아닐 경우 모두 빈 문자열로만 응답해줘
    """

    # 체인 생성
    if is_bytes:
        vision_chain = pass_image_chain | image_model | parser
    else:
        vision_chain = load_image_chain | image_model | parser

    return vision_chain.invoke({'image_path': image_path,
                                'prompt': vision_prompt},
                               )


# # # # # # # # # #
# 웹 페이지 구성
# # # # # # # # # #


st.set_page_config(layout="wide")
st.title('PPT 텍스트/이미지 분석봇')

# message key 등 세션 정보 초기화
if "messages" not in st.session_state:
    st.session_state.messages = []
if "api_key" not in st.session_state:  # api key
    st.session_state.api_key = None

# 이미지 파일 및 디스크립션들에 대한 캐시.
if 'img_showed_filenames' not in st.session_state:
    st.session_state.img_showed_filenames = []
if 'imgs' not in st.session_state:
    st.session_state.imgs = []
if 'descriptions' not in st.session_state:
    st.session_state.descriptions = []

# rag -> st.session_state.ppt_rag_chain
if 'ppt_docs' not in st.session_state:  # 모든 도큐먼트 목록. 파일이 추가될때마다 도큐먼트를 추가하고 새 rag를 생성한다.
    st.session_state.ppt_docs = []
    st.session_state.ppt_rag_chain = None
    st.session_state.ppt_vectorstore = None
if 'ppts_already_read_list' not in st.session_state:  # 이미 읽었던 ppt 파일 목록.
    st.session_state.ppts_already_read_list = []


# 사이드바 구성
if os.environ.get("OPENAI_API_KEY"):
    os.environ["OPENAI_API_KEY"] = st.sidebar.text_input('OpenAI API Key',
                                                         value=os.environ["OPENAI_API_KEY"])
else:
    os.environ["OPENAI_API_KEY"] = st.sidebar.text_input('OpenAI API Key',
                                                         placeholder='Input your ChatGPT API key here.')

# rag 관련 parameters
# max_input_len = st.sidebar.number_input('Max input length', min_value=1000, max_value=10000, value=5000, step=100)

user_files = st.sidebar.file_uploader('이미지 또는 PPT 파일을 올려주세요!', type=['jpg', 'jpeg', 'png', 'ppt', 'pptx'],
                                      accept_multiple_files=True)

# rag용 vector store 저장 버튼
if st.session_state.ppt_vectorstore:
    save_button = st.sidebar.button('도큐먼트 저장', type='primary')
    if save_button:
        st.session_state.ppt_vectorstore.save_local('./faiss_db')
        # 불러오기 코드
        # db_X = FAISS.load_local('./db/test_docs', embeddings=OpenAIEmbeddings(), allow_dangerous_deserialization=True)
else:
    load_button = st.sidebar.button('도큐먼트 불러오기', type='primary')
    if load_button:
        try:
            st.session_state.ppt_vectorstore = FAISS.load_local('./faiss_db', embeddings=OpenAIEmbeddings(),
                                                                allow_dangerous_deserialization=True)
        except:
            st.toast('DB 파일이 없습니다! (./faiss_db)', icon='🤬')

# 유저가 업로드한 파일 처리.
if user_files:
    for uploaded_file in user_files:
        # st.write("filename:", uploaded_file.name)
        logging.debug(str(uploaded_file))

        bytes_data = uploaded_file.read()

        # 기존에 출력했던 이미지일 경우.
        if uploaded_file.name in st.session_state.img_showed_filenames:
            ind = st.session_state.img_showed_filenames.index(uploaded_file.name)
            image = st.session_state.imgs[ind]
            desc = st.session_state.descriptions[ind]

            st.image(image, caption=uploaded_file.name, use_column_width=True)
            st.write(desc)

        # 이미지 파일일 경우 and 기존에 이미 출력한 파일이 아닐 경우에만.
        elif '.jpg' in uploaded_file.name or '.jpeg' in uploaded_file.name or '.png' in uploaded_file.name \
                and uploaded_file.name not in st.session_state.img_showed_filenames:
            image = Image.open(uploaded_file)
            # chat_placeholder.image(image, caption=uploaded_file.name)
            st.image(image, caption=uploaded_file.name, use_column_width=True)

            # logging.debug('image: ')
            # logging.debug(type(image))

            # 비동기 작업 제출
            future = executor.submit(get_image_description, bytes_data, True)

            # 작업 완료 여부 확인
            with st.spinner("이미지 디스크립션 작성 중..."):
                text_gen_result = future.result()  # 작업 완료 시까지 대기
            st.write(text_gen_result)

            # text_gen_result = get_image_description(bytes_data, is_bytes=True)
            # st.write(text_gen_result)

            st.session_state.img_showed_filenames.append(uploaded_file.name)
            st.session_state.imgs.append(image)
            st.session_state.descriptions.append(text_gen_result)

        # ppt 파일일 경우 and 기존에 읽었던 ppt 파일이 아닐 경우.
        elif '.ppt' in uploaded_file.name and uploaded_file.name not in st.session_state.ppts_already_read_list:
            # pprint(uploaded_file)
            # 새로운 파일을 읽은 파일 목록에 넣는다. (세션 유지 중 파일 다시 읽는것 방지)
            st.session_state.ppts_already_read_list.append(uploaded_file.name)

            # ppt 내 텍스트 doc 생성
            with st.spinner('파워포인트 파일 내 텍스트 읽는 중...'):
                docs1, tmp_pptxfile_path = get_text_documents_from_pptx(uploaded_file)

            st.write(docs1[:5])  # 텍스트 doc (일부만) 화면 출력

            with st.spinner('파워포인트 파일 내 이미지 읽는 중...'):
                # ppt 내 이미지를 모두 뽑아 doc 생성
                docs2, img_file_list = get_image_documents_from_pptx(tmp_pptxfile_path)

            # 이미지 & 이미지 doc 화면 출력
            for img_file, img_doc in zip(img_file_list, docs2):
                image = Image.open(img_file)
                # chat_placeholder.image(image, caption=uploaded_file.name)
                st.image(image, caption=img_file, use_column_width=True)
                st.write(img_doc)

            with st.spinner('RAG Chain 구성중...'):
                # 모든 도큐먼트를 합친다
                st.session_state.ppt_docs += combine_and_sort_documents(docs1, docs2)
                logging.debug('전체 docs 길이: ' + str(len(st.session_state.ppt_docs)))

                # rag를 만든다.
                st.session_state.ppt_rag_chain, \
                    st.session_state.ppt_vectorstore = get_rag_chain_from_docs(st.session_state.ppt_docs)
                st.write('파워포인트 내용 이해 완료! 무엇이 궁금하신가요?')


# 메인 페이지 구성
chat_placeholder = st.empty()

# 채팅 목록을 출력해줄 컨테이너 생성
text_container = st.container(border=True)
# text_container.title('AI와의 오붓한 채팅방')

for message in st.session_state.messages:
    with text_container:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

user_chat = st.chat_input("Say something")

if user_chat:
    st.session_state.messages.append({"role": "user", "content": user_chat})  # history 추가
    with text_container:
        with st.chat_message('user'):
            st.markdown(user_chat)

    # AI response
    if st.session_state.ppt_rag_chain:
        with text_container:
            with st.spinner('답변 생성중...'):
                answer = st.session_state.ppt_rag_chain.invoke(user_chat)

        st.session_state.messages.append({"role": "assistant", "content": answer})
        with text_container:
            with st.chat_message('assistant'):
                st.markdown(answer)
